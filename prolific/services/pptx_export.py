"""PowerPoint presentation generator for blog articles.

Two-phase approach:
1. LLM extracts content and writes speaker script → PresentationPlan
2. python-pptx renders slides with consistent dark theme
"""

import base64
import logging
import tempfile
from pathlib import Path

import httpx
from langchain_core.messages import HumanMessage, SystemMessage
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from prolific.schemas.presentation import PresentationPlan, SlideType
from prolific.services import pptx_theme as theme

logger = logging.getLogger(__name__)

SYSTEM_PROMPT = """\
You are a presentation designer creating a slide deck for a YouTube video.
Your job is to transform a technical blog article into a compelling, visually varied presentation.

SLIDE RULES:
- Cover ALL major sections of the article — do not skip any.
- Each KEY_POINTS slide should have 3-5 bullet points of concise phrases (NOT full sentences or paragraphs).
- Use a MIX of slide types: section_divider, key_points, image_feature, quote_highlight.
- Use quote_highlight for striking statistics, bold claims, or memorable insights.
- Use image_feature when a relevant image is available — reference images by their index number.
- Start with a TITLE slide and end with a CLOSING slide.
- Add SECTION_DIVIDER slides to separate major topic changes.
- Include a SOURCES slide near the end.
- Slide count should be proportional to article length: roughly 1 slide per 400-600 words of article text.

SPEAKER NOTES RULES:
- Speaker notes are a VERBATIM script for YouTube narration — the presenter will read them word-for-word.
- Write in a natural, conversational tone as if explaining to a knowledgeable audience.
- 3-5 sentences per slide.
- Start each slide's notes with a transition from the previous slide (e.g., "Now let's look at...", "Building on that...").
- Do NOT just repeat the bullet points — expand on them with context, examples, and insights from the article.
- For the TITLE slide, write an engaging opening hook that draws viewers in.
- For the CLOSING slide, write a strong summary and call-to-action (like, subscribe, check out the blog).

OUTPUT RULES:
- presentation_subtitle should be a one-line hook or tagline for the article.
- key_takeaway should be the single most important insight from the article.
- For IMAGE_FEATURE slides, set image_index to the index of the image from the AVAILABLE IMAGES list.
- If no images are listed, do not create image_feature slides.
"""


def _assemble_article_text(final_state: dict) -> tuple[str, str]:
    """Assemble article text from draft chunks, return (text, title)."""
    draft_chunks = final_state.get("draft_chunks", [])
    chapter_briefs = {b.chapter_id: b for b in final_state.get("chapter_briefs", [])}

    sorted_chunks = sorted(
        draft_chunks,
        key=lambda c: getattr(chapter_briefs.get(c.chapter_id), "chapter_number", 0),
    )

    parts = []
    for chunk in sorted_chunks:
        brief = chapter_briefs.get(chunk.chapter_id)
        if brief and len(sorted_chunks) > 1:
            parts.append(f"## {brief.title}\n\n{chunk.content}")
        else:
            parts.append(chunk.content)

    full_text = "\n\n".join(parts)

    global_memory = final_state.get("global_memory")
    title = getattr(global_memory, "title", "") if global_memory else ""
    return full_text, title


def _build_chapter_summary(final_state: dict) -> str:
    """Build a concise chapter structure summary for the LLM."""
    chapter_briefs = sorted(
        final_state.get("chapter_briefs", []),
        key=lambda b: b.chapter_number,
    )
    lines = []
    for b in chapter_briefs:
        points = ", ".join(b.key_points[:5]) if b.key_points else "N/A"
        lines.append(
            f"Chapter {b.chapter_number}: {b.title}\n"
            f"  Thesis: {b.thesis_statement}\n"
            f"  Key points: {points}"
        )
    return "\n\n".join(lines)


def _build_image_list(final_state: dict) -> list[dict]:
    """Build indexed list of available images with metadata."""
    visual_assets = final_state.get("visual_assets", [])
    images = []
    for asset in visual_assets:
        if asset.file_path or asset.url or asset.base64_data:
            images.append({
                "caption": asset.caption or asset.alt_text or "Image",
                "alt_text": asset.alt_text or "",
                "visual_type": asset.visual_type.value if hasattr(asset.visual_type, "value") else str(asset.visual_type),
                "asset": asset,
            })
    return images


def _build_claims_text(final_state: dict) -> str:
    """Extract high-confidence claims for quote slides."""
    claims = final_state.get("claims", [])
    high_claims = [c for c in claims if getattr(c, "confidence", None) and c.confidence.value == "high"]
    if not high_claims:
        high_claims = claims[:5]
    lines = [f"- {c.statement}" for c in high_claims[:10]]
    return "\n".join(lines) if lines else "No specific claims extracted."


def _build_sources_text(final_state: dict) -> str:
    """Build source summary for the LLM."""
    sources = final_state.get("approved_sources", [])
    if not sources:
        return "No sources available."
    lines = [f"- {s.title} ({s.url})" for s in sources[:8]]
    return f"Total sources: {len(sources)}\nTop sources:\n" + "\n".join(lines)


def _resolve_image_path(asset, blog_images_dir: Path) -> Path | None:
    """Resolve a VisualAsset to a local file path for embedding."""
    if asset.file_path:
        local_name = Path(asset.file_path).name
        blog_path = blog_images_dir / local_name
        if blog_path.exists():
            return blog_path
        original = Path(asset.file_path)
        if original.exists():
            return original

    if asset.url and not asset.url.startswith("data:"):
        try:
            resp = httpx.get(asset.url, timeout=15, follow_redirects=True)
            if resp.status_code == 200:
                content_type = resp.headers.get("content-type", "")
                if "png" in content_type:
                    suffix = ".png"
                elif "jpeg" in content_type or "jpg" in content_type:
                    suffix = ".jpg"
                elif "webp" in content_type:
                    suffix = ".webp"
                else:
                    suffix = Path(asset.url.split("?")[0]).suffix or ".jpg"
                tmp = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
                tmp.write(resp.content)
                tmp.close()
                return Path(tmp.name)
        except Exception as e:
            logger.warning(f"Failed to download image {asset.url}: {e}")

    if asset.base64_data:
        suffix = f".{asset.format}" if asset.format else ".png"
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
        tmp.write(base64.b64decode(asset.base64_data))
        tmp.close()
        return Path(tmp.name)

    return None


# -- Slide background helpers --

def _set_slide_bg(slide, color):
    """Set solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(slide, text, left, top, width, height,
                  font_size=Pt(18), font_color=theme.TEXT_PRIMARY,
                  bold=False, alignment=PP_ALIGN.LEFT,
                  font_name=theme.FONT_BODY, anchor=MSO_ANCHOR.TOP):
    """Add a styled text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_accent_line(slide, left, top, width):
    """Add a thin accent-colored line."""
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top, width, theme.ACCENT_LINE_HEIGHT,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme.ACCENT_PRIMARY
    shape.line.fill.background()
    return shape


def _add_speaker_notes(slide, notes_text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text


# -- Slide renderers --

def _render_title_slide(prs, slide_content, topic, date_str, presentation_title=None):
    """Render the opening title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    _set_slide_bg(slide, theme.BG_PRIMARY)

    _add_accent_line(slide, theme.TITLE_LEFT, Inches(1.8), theme.ACCENT_LINE_WIDTH)

    _add_text_box(
        slide, presentation_title or slide_content.title or topic,
        theme.TITLE_LEFT, Inches(2.1), theme.TITLE_WIDTH, Inches(2.0),
        font_size=theme.TITLE_FONT_SIZE, bold=True,
        font_color=theme.TEXT_PRIMARY, font_name=theme.FONT_TITLE,
    )

    if slide_content.subtitle:
        _add_text_box(
            slide, slide_content.subtitle,
            theme.TITLE_LEFT, Inches(4.2), theme.TITLE_WIDTH, Inches(1.0),
            font_size=theme.SUBTITLE_FONT_SIZE,
            font_color=theme.TEXT_SECONDARY,
        )

    _add_text_box(
        slide, f"{date_str}  |  Prolific Blog",
        theme.TITLE_LEFT, Inches(6.2), theme.TITLE_WIDTH, Inches(0.5),
        font_size=theme.DATE_FONT_SIZE,
        font_color=theme.TEXT_SECONDARY,
    )

    _add_speaker_notes(slide, slide_content.speaker_notes)


def _render_section_divider(slide_layout, prs, slide_content, section_num):
    """Render a section divider slide."""
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_bg(slide, theme.BG_SECTION)

    if section_num is not None:
        _add_text_box(
            slide, f"{section_num:02d}",
            theme.MARGIN_LEFT, Inches(1.5), Inches(2.0), Inches(1.2),
            font_size=theme.SECTION_NUMBER_SIZE, bold=True,
            font_color=theme.ACCENT_PRIMARY, font_name=theme.FONT_TITLE,
        )

    title_top = Inches(2.8) if section_num else Inches(2.5)
    _add_text_box(
        slide, slide_content.title,
        theme.MARGIN_LEFT, title_top, theme.CONTENT_WIDTH, Inches(1.5),
        font_size=theme.SECTION_TITLE_SIZE, bold=True,
        font_color=theme.TEXT_PRIMARY, font_name=theme.FONT_TITLE,
    )

    if slide_content.subtitle:
        _add_text_box(
            slide, slide_content.subtitle,
            theme.MARGIN_LEFT, title_top + Inches(1.6), theme.CONTENT_WIDTH, Inches(1.0),
            font_size=theme.SECTION_SUBTITLE_SIZE,
            font_color=theme.TEXT_SECONDARY,
        )

    _add_accent_line(slide, theme.MARGIN_LEFT, title_top - Inches(0.2), theme.ACCENT_LINE_WIDTH)
    _add_speaker_notes(slide, slide_content.speaker_notes)


def _render_key_points(slide_layout, prs, slide_content):
    """Render a bullet-point content slide."""
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_bg(slide, theme.BG_PRIMARY)

    _add_text_box(
        slide, slide_content.title,
        theme.MARGIN_LEFT, theme.MARGIN_TOP, theme.CONTENT_WIDTH, Inches(0.8),
        font_size=theme.SLIDE_TITLE_SIZE, bold=True,
        font_color=theme.TEXT_PRIMARY, font_name=theme.FONT_TITLE,
    )

    _add_accent_line(slide, theme.MARGIN_LEFT, theme.ACCENT_LINE_TOP, theme.ACCENT_LINE_WIDTH)

    bullet_top = theme.CONTENT_TOP
    for i, point in enumerate(slide_content.bullet_points[:6]):
        dot_color = theme.ACCENT_PRIMARY if i % 2 == 0 else theme.ACCENT_SECONDARY
        _add_text_box(
            slide, "\u2022",
            theme.MARGIN_LEFT, bullet_top, Inches(0.4), Inches(0.5),
            font_size=theme.BULLET_FONT_SIZE, font_color=dot_color, bold=True,
        )
        _add_text_box(
            slide, point,
            Inches(1.4), bullet_top, Inches(10.5), Inches(0.5),
            font_size=theme.BULLET_FONT_SIZE,
            font_color=theme.TEXT_PRIMARY,
        )
        bullet_top += Inches(0.75)

    _add_speaker_notes(slide, slide_content.speaker_notes)


def _render_image_feature(slide_layout, prs, slide_content, image_path):
    """Render an image-focused slide with optional caption and bullets."""
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_bg(slide, theme.BG_PRIMARY)

    _add_text_box(
        slide, slide_content.title,
        theme.MARGIN_LEFT, theme.MARGIN_TOP, theme.CONTENT_WIDTH, Inches(0.8),
        font_size=theme.SLIDE_TITLE_SIZE, bold=True,
        font_color=theme.TEXT_PRIMARY, font_name=theme.FONT_TITLE,
    )

    _add_accent_line(slide, theme.MARGIN_LEFT, theme.ACCENT_LINE_TOP, theme.ACCENT_LINE_WIDTH)

    if image_path and image_path.exists():
        try:
            slide.shapes.add_picture(
                str(image_path),
                theme.IMAGE_LEFT, theme.IMAGE_TOP,
                theme.IMAGE_WIDTH, theme.IMAGE_HEIGHT,
            )
        except Exception as e:
            logger.warning(f"Failed to embed image {image_path}: {e}")
            _add_text_box(
                slide, "[Image unavailable]",
                theme.IMAGE_LEFT, Inches(3.5), theme.IMAGE_WIDTH, Inches(1.0),
                font_size=theme.BULLET_FONT_SIZE,
                font_color=theme.TEXT_SECONDARY,
                alignment=PP_ALIGN.CENTER,
            )

    caption = slide_content.image_caption or ""
    if caption:
        _add_text_box(
            slide, caption,
            theme.IMAGE_LEFT, Inches(6.5), theme.IMAGE_WIDTH, Inches(0.5),
            font_size=theme.CAPTION_FONT_SIZE,
            font_color=theme.TEXT_SECONDARY,
            alignment=PP_ALIGN.CENTER,
        )

    if slide_content.bullet_points:
        text_top = theme.IMAGE_TOP
        for point in slide_content.bullet_points[:4]:
            _add_text_box(
                slide, f"\u2022  {point}",
                theme.IMAGE_TEXT_LEFT, text_top, theme.IMAGE_TEXT_WIDTH, Inches(0.6),
                font_size=theme.BULLET_SUB_SIZE,
                font_color=theme.TEXT_PRIMARY,
            )
            text_top += Inches(0.7)

    _add_speaker_notes(slide, slide_content.speaker_notes)


def _render_quote_highlight(slide_layout, prs, slide_content):
    """Render a large quote/stat highlight slide."""
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_bg(slide, theme.BG_PRIMARY)

    quote_bg = slide.shapes.add_shape(
        1, Inches(1.5), Inches(1.5), Inches(10.3), Inches(4.5),
    )
    quote_bg.fill.solid()
    quote_bg.fill.fore_color.rgb = theme.HIGHLIGHT_BG
    quote_bg.line.fill.background()

    _add_text_box(
        slide, "\u201C",
        Inches(1.8), Inches(1.3), Inches(1.0), Inches(1.0),
        font_size=Pt(72), font_color=theme.QUOTE_MARK_COLOR,
        bold=True, font_name=theme.FONT_TITLE,
    )

    quote = slide_content.quote_text or slide_content.title
    _add_text_box(
        slide, quote,
        Inches(2.0), Inches(2.3), Inches(9.3), Inches(2.5),
        font_size=theme.QUOTE_FONT_SIZE,
        font_color=theme.TEXT_PRIMARY,
        alignment=PP_ALIGN.LEFT,
    )

    if slide_content.quote_attribution:
        _add_text_box(
            slide, f"\u2014 {slide_content.quote_attribution}",
            Inches(2.0), Inches(5.0), Inches(9.3), Inches(0.5),
            font_size=theme.ATTRIBUTION_FONT_SIZE,
            font_color=theme.TEXT_SECONDARY,
            alignment=PP_ALIGN.RIGHT,
        )

    _add_speaker_notes(slide, slide_content.speaker_notes)


def _render_comparison(slide_layout, prs, slide_content):
    """Render a comparison slide with side-by-side columns."""
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_bg(slide, theme.BG_PRIMARY)

    _add_text_box(
        slide, slide_content.title,
        theme.MARGIN_LEFT, theme.MARGIN_TOP, theme.CONTENT_WIDTH, Inches(0.8),
        font_size=theme.SLIDE_TITLE_SIZE, bold=True,
        font_color=theme.TEXT_PRIMARY, font_name=theme.FONT_TITLE,
    )

    _add_accent_line(slide, theme.MARGIN_LEFT, theme.ACCENT_LINE_TOP, theme.ACCENT_LINE_WIDTH)

    points = slide_content.bullet_points[:6]
    mid = len(points) // 2 or 1
    left_points = points[:mid]
    right_points = points[mid:]

    col_top = theme.CONTENT_TOP
    for point in left_points:
        _add_text_box(
            slide, f"\u2022  {point}",
            Inches(0.8), col_top, Inches(5.5), Inches(0.6),
            font_size=theme.BULLET_FONT_SIZE,
            font_color=theme.TEXT_PRIMARY,
        )
        col_top += Inches(0.75)

    col_top = theme.CONTENT_TOP
    for point in right_points:
        _add_text_box(
            slide, f"\u2022  {point}",
            Inches(7.0), col_top, Inches(5.5), Inches(0.6),
            font_size=theme.BULLET_FONT_SIZE,
            font_color=theme.TEXT_PRIMARY,
        )
        col_top += Inches(0.75)

    _add_speaker_notes(slide, slide_content.speaker_notes)


def _render_sources(slide_layout, prs, slide_content, approved_sources):
    """Render a sources/references slide."""
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_bg(slide, theme.BG_PRIMARY)

    _add_text_box(
        slide, "Sources & References",
        theme.MARGIN_LEFT, theme.MARGIN_TOP, theme.CONTENT_WIDTH, Inches(0.8),
        font_size=theme.SLIDE_TITLE_SIZE, bold=True,
        font_color=theme.TEXT_PRIMARY, font_name=theme.FONT_TITLE,
    )

    _add_accent_line(slide, theme.MARGIN_LEFT, theme.ACCENT_LINE_TOP, theme.ACCENT_LINE_WIDTH)

    sources_to_show = approved_sources[:10]
    src_top = theme.CONTENT_TOP
    for i, src in enumerate(sources_to_show):
        col_left = Inches(0.8) if i < 5 else Inches(7.0)
        row_top = src_top + Inches((i % 5) * 0.85)
        title = src.title[:70] + "..." if len(src.title) > 70 else src.title
        _add_text_box(
            slide, f"{i + 1}. {title}",
            col_left, row_top, Inches(5.5), Inches(0.4),
            font_size=theme.SOURCE_FONT_SIZE,
            font_color=theme.TEXT_PRIMARY, bold=True,
        )
        url_display = src.url[:60] + "..." if len(src.url) > 60 else src.url
        _add_text_box(
            slide, url_display,
            col_left, row_top + Inches(0.35), Inches(5.5), Inches(0.3),
            font_size=Pt(9),
            font_color=theme.TEXT_SECONDARY,
        )

    _add_text_box(
        slide, "Full article with all citations available on the blog",
        theme.MARGIN_LEFT, Inches(6.5), theme.CONTENT_WIDTH, Inches(0.5),
        font_size=theme.CAPTION_FONT_SIZE,
        font_color=theme.ACCENT_PRIMARY,
        alignment=PP_ALIGN.CENTER,
    )

    _add_speaker_notes(slide, slide_content.speaker_notes)


def _render_closing(slide_layout, prs, slide_content, key_takeaway):
    """Render the closing slide with key takeaway and CTA."""
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_bg(slide, theme.BG_SECTION)

    _add_accent_line(slide, theme.TITLE_LEFT, Inches(2.0), theme.ACCENT_LINE_WIDTH)

    takeaway = key_takeaway or slide_content.title
    _add_text_box(
        slide, takeaway,
        theme.TITLE_LEFT, Inches(2.3), theme.TITLE_WIDTH, Inches(2.0),
        font_size=Pt(30), bold=True,
        font_color=theme.TEXT_PRIMARY, font_name=theme.FONT_TITLE,
        alignment=PP_ALIGN.LEFT,
    )

    if slide_content.subtitle:
        _add_text_box(
            slide, slide_content.subtitle,
            theme.TITLE_LEFT, Inches(4.5), theme.TITLE_WIDTH, Inches(1.0),
            font_size=theme.SUBTITLE_FONT_SIZE,
            font_color=theme.TEXT_SECONDARY,
        )

    _add_text_box(
        slide, "Like & Subscribe  |  Prolific Blog",
        theme.TITLE_LEFT, Inches(6.2), theme.TITLE_WIDTH, Inches(0.5),
        font_size=theme.DATE_FONT_SIZE,
        font_color=theme.ACCENT_PRIMARY,
        alignment=PP_ALIGN.LEFT,
    )

    _add_speaker_notes(slide, slide_content.speaker_notes)


# -- Result type --

class PresentationResult:
    """Result of presentation generation with metrics for monitoring."""

    def __init__(self):
        self.status: str = "pending"
        self.file_path: str | None = None
        self.slide_count: int = 0
        self.images_available: int = 0
        self.images_embedded: int = 0
        self.images_failed: int = 0
        self.slide_types: dict[str, int] = {}
        self.has_speaker_notes: bool = False
        self.duration_seconds: float = 0.0
        self.error: str | None = None
        self.error_stage: str | None = None

    def to_dict(self) -> dict:
        return {
            "status": self.status,
            "file_path": self.file_path,
            "slide_count": self.slide_count,
            "images_available": self.images_available,
            "images_embedded": self.images_embedded,
            "images_failed": self.images_failed,
            "slide_types": self.slide_types,
            "has_speaker_notes": self.has_speaker_notes,
            "duration_seconds": round(self.duration_seconds, 1),
            "error": self.error,
            "error_stage": self.error_stage,
        }


# -- Main generation function --

async def generate_presentation(
    final_state: dict,
    slug: str,
    topic: str,
    project_root: Path,
    blog_images_dir: Path,
) -> PresentationResult:
    """Generate a PowerPoint presentation from article generation state.

    Returns a PresentationResult with metrics, status, and error details.
    """
    import time
    start = time.monotonic()
    result = PresentationResult()

    # -- Stage 1: Validate inputs --
    draft_chunks = final_state.get("draft_chunks", [])
    if not draft_chunks:
        result.status = "skipped"
        result.error = "No draft chunks available"
        result.error_stage = "validation"
        result.duration_seconds = time.monotonic() - start
        logger.warning("PPTX skipped: no draft chunks")
        return result

    logger.info("=== PRESENTATION GENERATION ===")

    # -- Stage 2: Assemble article data --
    try:
        article_text, title = _assemble_article_text(final_state)
        if not title:
            title = topic

        chapter_summary = _build_chapter_summary(final_state)
        image_list = _build_image_list(final_state)
        claims_text = _build_claims_text(final_state)
        sources_text = _build_sources_text(final_state)

        result.images_available = len(image_list)
        word_count = sum(getattr(c, "word_count", 0) for c in draft_chunks)

        logger.info(f"[1/4] Data assembled: {word_count} words, "
                     f"{len(image_list)} images, "
                     f"{len(final_state.get('chapter_briefs', []))} chapters")
    except Exception as e:
        result.status = "failed"
        result.error = f"Data assembly failed: {e}"
        result.error_stage = "data_assembly"
        result.duration_seconds = time.monotonic() - start
        logger.error(f"PPTX failed at data assembly: {e}", exc_info=True)
        return result

    # -- Stage 3: LLM content extraction --
    try:
        image_list_str = ""
        if image_list:
            lines = []
            for i, img in enumerate(image_list):
                lines.append(f"[{i}] Type: {img['visual_type']}, Caption: {img['caption']}, Alt: {img['alt_text']}")
            image_list_str = "\n".join(lines)
        else:
            image_list_str = "No images available. Do not create image_feature slides."

        max_article_chars = 30000
        if len(article_text) > max_article_chars:
            logger.info(f"Article text truncated: {len(article_text)} -> {max_article_chars} chars")
            article_text = article_text[:max_article_chars] + "\n\n[Article truncated for length...]"

        from datetime import datetime
        date_str = datetime.now().strftime("%B %d, %Y")

        human_msg = (
            f"ARTICLE TITLE: {title}\n"
            f"DATE: {date_str}\n"
            f"WORD COUNT: {word_count}\n\n"
            f"CHAPTER STRUCTURE:\n{chapter_summary}\n\n"
            f"AVAILABLE IMAGES (reference by index number):\n{image_list_str}\n\n"
            f"HIGH-CONFIDENCE CLAIMS (use for quote slides):\n{claims_text}\n\n"
            f"SOURCES:\n{sources_text}\n\n"
            f"FULL ARTICLE TEXT:\n{article_text}\n\n"
            f"Create a presentation plan for this article."
        )

        logger.info("[2/4] Calling LLM for presentation plan...")
        from prolific.services.llm import get_llm_service
        llm_service = get_llm_service()

        plan: PresentationPlan = await llm_service.invoke_with_structured_output(
            messages=[
                SystemMessage(content=SYSTEM_PROMPT),
                HumanMessage(content=human_msg),
            ],
            output_schema=PresentationPlan,
            tier="extraction",
            temperature=0.4,
        )

        result.slide_count = len(plan.slides)
        type_counts: dict[str, int] = {}
        for s in plan.slides:
            t = s.slide_type.value
            type_counts[t] = type_counts.get(t, 0) + 1
        result.slide_types = type_counts
        result.has_speaker_notes = all(s.speaker_notes.strip() for s in plan.slides)

        logger.info(f"[2/4] LLM plan received: {result.slide_count} slides, "
                     f"types={type_counts}")

    except Exception as e:
        result.status = "failed"
        result.error = f"LLM planning failed: {e}"
        result.error_stage = "llm_planning"
        result.duration_seconds = time.monotonic() - start
        logger.error(f"PPTX failed at LLM planning: {e}", exc_info=True)
        return result

    # -- Stage 4: Resolve images --
    resolved_images: dict[int, Path | None] = {}
    temp_files: list[Path] = []
    try:
        for i, img in enumerate(image_list):
            path = _resolve_image_path(img["asset"], blog_images_dir)
            resolved_images[i] = path
            if path and str(path).startswith(tempfile.gettempdir()):
                temp_files.append(path)
            if path:
                logger.info(f"[3/4] Image [{i}] resolved: {path.name}")
            else:
                logger.warning(f"[3/4] Image [{i}] could not be resolved")

        resolved_count = sum(1 for p in resolved_images.values() if p is not None)
        logger.info(f"[3/4] Images resolved: {resolved_count}/{len(image_list)}")
    except Exception as e:
        logger.warning(f"[3/4] Image resolution error (non-fatal): {e}")

    # -- Stage 5: Render PPTX --
    try:
        prs = Presentation()
        prs.slide_width = theme.SLIDE_WIDTH
        prs.slide_height = theme.SLIDE_HEIGHT

        blank_layout = prs.slide_layouts[6]
        section_counter = 0
        images_embedded = 0
        images_failed = 0
        approved_sources = final_state.get("approved_sources", [])

        for idx, slide_content in enumerate(plan.slides):
            st = slide_content.slide_type

            if st == SlideType.TITLE:
                _render_title_slide(prs, slide_content, title, date_str,
                                    presentation_title=plan.presentation_title)

            elif st == SlideType.SECTION_DIVIDER:
                section_counter += 1
                sec_num = slide_content.section_number or section_counter
                _render_section_divider(blank_layout, prs, slide_content, sec_num)

            elif st == SlideType.KEY_POINTS:
                _render_key_points(blank_layout, prs, slide_content)

            elif st == SlideType.IMAGE_FEATURE:
                img_path = None
                if slide_content.image_index is not None:
                    img_path = resolved_images.get(slide_content.image_index)
                if img_path and img_path.exists():
                    _render_image_feature(blank_layout, prs, slide_content, img_path)
                    images_embedded += 1
                else:
                    if slide_content.image_index is not None:
                        images_failed += 1
                        logger.warning(f"[4/4] Slide {idx}: image [{slide_content.image_index}] "
                                       f"unavailable, falling back to key_points")
                    _render_key_points(blank_layout, prs, slide_content)

            elif st == SlideType.QUOTE_HIGHLIGHT:
                _render_quote_highlight(blank_layout, prs, slide_content)

            elif st == SlideType.COMPARISON:
                _render_comparison(blank_layout, prs, slide_content)

            elif st == SlideType.SOURCES:
                _render_sources(blank_layout, prs, slide_content, approved_sources)

            elif st == SlideType.CLOSING:
                _render_closing(blank_layout, prs, slide_content, plan.key_takeaway)

        result.images_embedded = images_embedded
        result.images_failed = images_failed

        # Save
        output_dir = project_root / "blog" / "public" / "presentations"
        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_dir / f"{slug}.pptx"
        prs.save(str(output_path))

        result.status = "success"
        result.file_path = str(output_path)
        result.duration_seconds = time.monotonic() - start

        logger.info(f"[4/4] Presentation saved: {output_path}")
        logger.info(f"=== PRESENTATION COMPLETE: {result.slide_count} slides, "
                     f"{images_embedded} images embedded, "
                     f"{result.duration_seconds:.1f}s ===")
        return result

    except Exception as e:
        result.status = "failed"
        result.error = f"PPTX rendering failed: {e}"
        result.error_stage = "rendering"
        result.duration_seconds = time.monotonic() - start
        logger.error(f"PPTX failed at rendering: {e}", exc_info=True)
        return result

    finally:
        for tmp in temp_files:
            try:
                tmp.unlink(missing_ok=True)
            except Exception:
                pass
